"""Build a leadership briefing deck for mRNA-EditFlow.

The deck is intentionally evidence-driven: most headline numbers are loaded
from current benchmark/readiness artifacts, and claim boundaries are kept
explicit so the slides can be used for management reporting without
over-claiming.
"""
from __future__ import annotations

import json
import math
import os
from datetime import date
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
OUT_PPTX = OUT_DIR / "mrna_editflow_leadership_briefing_20260715.pptx"
OUT_MD = OUT_DIR / "mrna_editflow_leadership_briefing_20260715_outline.md"

COLORS = {
    "navy": RGBColor(18, 43, 84),
    "blue": RGBColor(42, 100, 180),
    "cyan": RGBColor(32, 150, 180),
    "green": RGBColor(42, 157, 92),
    "orange": RGBColor(224, 132, 47),
    "red": RGBColor(190, 70, 70),
    "purple": RGBColor(111, 86, 180),
    "gray": RGBColor(95, 105, 120),
    "light_gray": RGBColor(238, 242, 247),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(30, 35, 45),
}


def load_json(rel: str) -> dict[str, Any]:
    path = ROOT / rel
    if not path.exists():
        return {}
    with open(path, "r", encoding="utf-8") as fh:
        payload = json.load(fh)
    return payload if isinstance(payload, dict) else {}


def agg_mean(rel: str, key: str) -> float | None:
    payload = load_json(rel)
    entry = payload.get("aggregate", {}).get(key)
    if isinstance(entry, Mapping):
        value = entry.get("mean")
        return float(value) if isinstance(value, (int, float)) else None
    return None


def fmt(value: Any, digits: int = 4, signed: bool = False) -> str:
    if value is None:
        return "NA"
    if isinstance(value, bool):
        return "True" if value else "False"
    try:
        x = float(value)
    except (TypeError, ValueError):
        return str(value)
    if not math.isfinite(x):
        return "NA"
    sign = "+" if signed and x >= 0 else ""
    return f"{sign}{x:.{digits}f}"


def pct(value: Any, digits: int = 1) -> str:
    if value is None:
        return "NA"
    try:
        return f"{float(value) * 100:.{digits}f}%"
    except (TypeError, ValueError):
        return str(value)


def set_text(
    shape,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = "PingFang SC"
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.58)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["navy"]
    bar.line.color.rgb = COLORS["navy"]
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(12.4), Inches(0.35))
    set_text(title_box, title, size=20, bold=True, color=COLORS["white"])
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(12.0), Inches(0.3))
        set_text(sub, subtitle, size=11, color=COLORS["gray"])


def add_footer(slide, idx: int):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(7.16), Inches(12.45), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 225, 233)
    line.line.color.rgb = RGBColor(220, 225, 233)
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(7.19), Inches(10.8), Inches(0.25))
    set_text(
        foot,
        "mRNA-EditFlow | proxy/offline evidence; no wet-lab or external SOTA claim unless gates are complete",
        size=8,
        color=COLORS["gray"],
    )
    page = slide.shapes.add_textbox(Inches(12.1), Inches(7.19), Inches(0.7), Inches(0.25))
    set_text(page, str(idx), size=8, color=COLORS["gray"], align=PP_ALIGN.RIGHT)


def add_bullets(
    slide,
    bullets: Sequence[str],
    *,
    x: float = 0.72,
    y: float = 1.25,
    w: float = 12.0,
    h: float = 5.65,
    size: int = 16,
    color: RGBColor | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = "PingFang SC"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["black"]
    return box


def add_metric_cards(slide, cards: Sequence[tuple[str, str, str]], *, y: float = 1.35):
    n = len(cards)
    gap = 0.18
    width = (12.15 - gap * (n - 1)) / n
    for i, (label, value, note) in enumerate(cards):
        x = 0.6 + i * (width + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(1.1),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(246, 249, 252)
        card.line.color.rgb = RGBColor(210, 218, 229)
        box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.10), Inches(width - 0.24), Inches(0.25))
        set_text(box, label, size=9, bold=True, color=COLORS["gray"])
        val = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.38), Inches(width - 0.24), Inches(0.32))
        set_text(val, value, size=17, bold=True, color=COLORS["blue"])
        nt = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.78), Inches(width - 0.24), Inches(0.25))
        set_text(nt, note, size=8, color=COLORS["gray"])


def add_table(slide, headers: Sequence[str], rows: Sequence[Sequence[str]], *, x=0.55, y=1.22, w=12.2, h=5.5, font=8):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]
        for p in cell.text_frame.paragraphs:
            p.font.name = "PingFang SC"
            p.font.bold = True
            p.font.size = Pt(font)
            p.font.color.rgb = COLORS["white"]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 252, 255) if r % 2 else RGBColor(240, 245, 250)
            for p in cell.text_frame.paragraphs:
                p.font.name = "PingFang SC"
                p.font.size = Pt(font)
                p.font.color.rgb = COLORS["black"]
    return table_shape


def add_bar_chart(
    slide,
    labels: Sequence[str],
    values: Sequence[float],
    *,
    x=0.9,
    y=1.45,
    w=11.6,
    h=3.6,
    title: str | None = None,
    color: RGBColor = COLORS["blue"],
):
    if title:
        box = slide.shapes.add_textbox(Inches(x), Inches(y - 0.45), Inches(w), Inches(0.3))
        set_text(box, title, size=12, bold=True, color=COLORS["black"])
    max_abs = max(abs(v) for v in values) or 1.0
    max_v = max(values)
    min_v = min(values)
    baseline = y + h if min_v >= 0 else y + h * (max_v / (max_v - min_v))
    axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(baseline), Inches(w), Inches(0.015))
    axis.fill.solid()
    axis.fill.fore_color.rgb = RGBColor(190, 198, 210)
    axis.line.color.rgb = RGBColor(190, 198, 210)
    bw = w / len(values) * 0.62
    spacing = w / len(values)
    for i, (label, val) in enumerate(zip(labels, values)):
        bx = x + i * spacing + (spacing - bw) / 2
        bh = h * abs(val) / max_abs * 0.85
        by = baseline - bh if val >= 0 else baseline
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(bx), Inches(by), Inches(bw), Inches(bh))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color if val >= 0 else COLORS["red"]
        rect.line.color.rgb = rect.fill.fore_color.rgb
        valbox = slide.shapes.add_textbox(Inches(bx - 0.05), Inches(by - 0.25 if val >= 0 else by + bh + 0.02), Inches(bw + 0.1), Inches(0.22))
        set_text(valbox, fmt(val, 4, signed=True), size=8, color=COLORS["black"], align=PP_ALIGN.CENTER)
        lab = slide.shapes.add_textbox(Inches(bx - 0.12), Inches(y + h + 0.08), Inches(bw + 0.24), Inches(0.42))
        set_text(lab, label, size=7, color=COLORS["gray"], align=PP_ALIGN.CENTER)


def slide(prs: Presentation, title: str, subtitle: str | None = None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, title, subtitle)
    return s


def build_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides_meta: list[tuple[str, list[str]]] = []

    readiness = load_json("docs/sota_readiness_audit_head256.json").get("summary", {})
    data_ready = load_json("docs/data_scaleup_readiness.json").get("summary", {})
    remote = load_json("docs/remote_execution_status.json").get("summary", {})
    t4_1024 = load_json("benchmark/t4_protein_identity_cai_gc_report_head1024.json")
    t4_metrics = load_json("benchmark/protein_conditioned_codon_metrics_head1024.json").get("summary", {})
    t4_cds = load_json("benchmark/protein_conditioned_cds_head1024.summary.json").get("summary", {})
    t4_sweep = load_json("benchmark/protein_conditioned_cds_gc_sweep_head1024.summary.json")
    t4_sweep_audit = load_json("benchmark/protein_conditioned_cds_gc_sweep_head1024.audit.json").get("summary", {})
    spectrum = load_json("benchmark/multi_scale_sequence_spectrum_head32_ranker_full1k.json").get("summary", {})

    t5_head256 = {
        "grpo": agg_mean("benchmark/multiseed_t5_public_head256_mo_grpo_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
        "scalar": agg_mean("benchmark/multiseed_t5_public_head256_mo_scalar_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
        "pareto": agg_mean("benchmark/multiseed_t5_public_head256_mo_pareto_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
        "te_only": agg_mean("benchmark/multiseed_t5_public_head256_mo_te_only_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
    }
    t5_head1024 = {
        "pareto": agg_mean("benchmark/multiseed_t5_public_head1024_mo_pareto_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
        "te_only": agg_mean("benchmark/multiseed_t5_public_head1024_mo_te_only_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
        "hardneg": agg_mean("benchmark/multiseed_t5_public_head1024_hardneg_v2_top64/multiseed_summary.json", "delta_oracle_te_vs_source"),
    }

    # 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.color.rgb = COLORS["navy"]
    title_box = s.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.8), Inches(1.1))
    set_text(title_box, "mRNA-EditFlow 项目汇报", size=38, bold=True, color=COLORS["white"])
    sub = s.shapes.add_textbox(Inches(0.85), Inches(2.45), Inches(11.5), Inches(1.0))
    set_text(
        sub,
        "全长 mRNA-native constrained edit-flow：数据、模型、评测、当前结果与下一步规划",
        size=20,
        color=RGBColor(218, 230, 245),
    )
    meta = s.shapes.add_textbox(Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.5))
    set_text(meta, f"汇报日期：{date.today().isoformat()} | 证据来源：当前仓库 JSON/MD artifacts", size=12, color=RGBColor(218, 230, 245))
    slides_meta.append(("封面", ["mRNA-EditFlow 项目汇报"]))

    # 2 executive summary
    s = slide(prs, "一页摘要：当前项目状态")
    add_metric_cards(
        s,
        [
            ("内部证据完整度", "6/6 ready", "SOTA readiness sections"),
            ("正向 SOTA claim", "False", "外部/湿实验/scale-law gate 未完成"),
            ("T1-T7 bundle", "9/9 reports", "proxy + hard-constraint evidence"),
            ("远端 artifact", f"{remote.get('n_artifacts_present', 0)} present", "remote status latest snapshot"),
        ],
    )
    add_bullets(
        s,
        [
            "已能做约束安全的 full-transcript local optimization / reranking：legal/protein/frame/budget 硬约束 exact-1。",
            "head256 多目标 ranker 严格显著优于同 recipe TE-only；head1024 上收益收缩，pareto 仅为 borderline/trend。",
            "protein-conditioned T4 已从 head256 扩到 head1024，codon-level audit ready，但仍不是 codonGPT/LinearDesign 外部复现。",
            "核心短板：真实 MPRA/stability 数据、外部 executable baseline、真实 frozen FM embedding、RefSeq/cross-family data diversity gate。",
        ],
        y=2.75,
        size=15,
    )
    slides_meta.append(("一页摘要", ["内部 proxy 证据 ready，但 positive SOTA claim 仍关闭"]))

    # 3 background
    s = slide(prs, "项目背景：为什么需要 mRNA-native edit-flow？")
    add_bullets(
        s,
        [
            "mRNA 不是普通字符串：5'UTR / CDS / 3'UTR 三段语法、功能和约束完全不同。",
            "CDS 必须保持 reading frame、起始/终止密码子、无 premature stop、protein identity exact-1。",
            "UTR 是调控画布：Kozak、uAUG、polyA、ARE、局部结构可及性影响 TE/stability。",
            "治疗性 construct 更常见需求是“在已有序列上最小编辑优化”，而不是自由重生成。",
            "MEF 的目标：在硬约束下，同时优化 UTR regulatory signal、CDS codon/structure property、full-transcript distribution。",
        ],
    )
    slides_meta.append(("项目背景", ["mRNA-native 三段语法和最小编辑需求"]))

    # 4 research questions
    s = slide(prs, "论文级核心问题")
    add_table(
        s,
        ["问题", "需要回答的证据"],
        [
            ["Full transcript design 是否更强？", "T1-T7 full transcript vs CDS-only / UTR-only / representation-only"],
            ["Constrained edit-flow 是否稳定提升？", "TE/stability proxy uplift + hard constraints exact-1"],
            ["Cascade decoding 是否提升 recall/top-1？", "candidate recall、top-1 design quality、paired tests"],
            ["Frozen mRNA FM 是否提升 ranking/generalization？", "leakage-free split + matched trainable budget adapter probe"],
            ["Scale data/params 是否规律提升？", "data × model × steps + diversity/complexity axes"],
            ["如何避免只优化低阶 proxy？", "multi-scale spectrum、perceptual/frozen-feature、adversarial realism audit"],
        ],
        h=4.85,
        font=10,
    )
    slides_meta.append(("论文问题", ["6 个 paper-level questions"]))

    # 5 one-sentence contribution
    s = slide(prs, "项目定位：MEF 和现有 SOTA 的差异")
    add_table(
        s,
        ["竞品类别", "代表方法", "MEF 差异化"],
        [
            ["CDS optimizer", "LinearDesign / EnsembleDesign", "MEF 保留 CDS 同义格，但扩展到 full transcript + edit budget"],
            ["5'UTR generator", "UTailoR / UTRGAN", "MEF 支持局部插删/模板改造，不只 de novo UTR"],
            ["Full-length AR", "GEMORNA / mRNA-GPT", "MEF 强调 constrained local edit-flow + 可解释 edit distance"],
            ["Protein-conditioned MO", "ProMORNA", "借鉴 per-metric advantage 标准化；MEF 需对齐 Pareto frontier"],
            ["Latent optimization", "RNAGenScape", "共同点是 optimize-not-invent；MEF 是离散可审计编辑轨迹"],
            ["Frozen FM", "mRNA-LM / Helix / CodonFM", "MEF 使用 frozen features 作为 edit-rate 条件，而不是仅做 prediction"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("SOTA 差异化", ["优化而非重造、硬约束、可解释 edit-flow"]))

    # 6 data overview
    s = slide(prs, "预训练与训练数据：当前资产和目标资产")
    add_table(
        s,
        ["数据源", "当前状态", "用途", "风险/缺口"],
        [
            ["GENCODE v45 human protein-coding", "54,680 cleaned records", "Stage A 主训练；T1/T4/T5/T6/T7", "human-only；仍需 cross-source panel"],
            ["RefSeq human RNA", "build running，records/manifest 未落盘", "data scale-up + GENCODE→RefSeq 泛化", "官方 raw/records/manifest 缺失"],
            ["MPRA/TE", "协议/manifest builder ready，真实表未接入", "真实 TE predictor / oracle", "不能 claim real TE"],
            ["Stability/half-life", "协议 ready，真实标签未接入", "stability downstream probe", "不能 claim real stability"],
            ["mRNA-LM/Helix/CodonFM", "文献/协议调研，真实 embedding cache 未接入", "frozen FM feature transfer", "leakage & license gate"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("数据资产", ["GENCODE 已可用；RefSeq/MPRA/stability/FMs 仍缺"]))

    # 7 dataset acquisition
    s = slide(prs, "数据集获取方法：从 public source 到 canonical records")
    add_bullets(
        s,
        [
            "GENCODE：固定 release，下载 transcript FASTA，解析 CDS range，抽取 5'UTR/CDS/3'UTR。",
            "RefSeq：保守 GenBank parser；只接收 plus-strand contiguous / simple join CDS，跳过 complement、remote accession、多 CDS。",
            "MPRA/Stability：要求真实 CSV/TSV + source URL + official split；缺官方 split 时拒绝生成 paper-grade manifest。",
            "所有版本记录 source URL、raw SHA256、clean SHA256、record count、drop stats、split sidecar。",
            "远端队列：RefSeq public build、GENCODE/RefSeq family leakage、downstream predictor protocol、P3 readiness watcher。",
        ],
    )
    slides_meta.append(("数据获取", ["public pipeline + conservative parser + manifest contract"]))

    # 8 data cleaning
    s = slide(prs, "数据处理方法：清洗、分区、manifest 与 leakage audit")
    add_table(
        s,
        ["处理步骤", "规则", "当前证据"],
        [
            ["序列规范化", "T→U；只允许 A/C/G/U；剔除 N/非法字符", "GENCODE records SHA verified"],
            ["CDS 合法性", "AUG start；terminal stop；len%3=0；no internal stop", "cleaned kept=54,680"],
            ["区域边界", "5'UTR / CDS / 3'UTR 明确；超长按配置截断", "stage_a_full_a100_max config"],
            ["family split", "gene/protein family + k-mer leakage sidecar", "GENCODE split ready；RefSeq waiting"],
            ["manifest audit", "source URL、SHA、counts、drop stats、split stats", "dataset contract not complete"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("数据处理", ["canonical records + manifest + leakage"]))

    # 9 eFold
    s = slide(prs, "eFold 启发：data scale-up 不只看条数")
    add_bullets(
        s,
        [
            "eFold / RNAndria 结论：merely expanding database size is insufficient for generalization across families；多样性和复杂度更关键。",
            "迁移到 MEF：data scale-up 从 data_size × model_size × steps 扩展为 family_diversity × region/length_complexity × motif/structure_complexity。",
            "新增 cross-family stress layer：gene/protein-family、motif-family、length/complexity、cross-source、cross-species。",
            "新增 gate：ready_for_cross_family_generalization_claim 只有在 diversity profile、complexity bucket、cross-source panel 都完成后才能打开。",
            "边界：eFold 是 RNA structure prediction，不是 mRNA design baseline；只迁移数据治理与泛化评测原则。",
        ],
    )
    slides_meta.append(("eFold 启发", ["数据多样性/复杂度和 cross-family gate"]))

    # 10 algorithm principle
    s = slide(prs, "算法原理：全长 mRNA-native constrained edit-flow")
    add_bullets(
        s,
        [
            "状态表示：x = (5'UTR, CDS, 3'UTR)，每个位置带 region id、codon phase、mask、source context。",
            "编辑算子：UTR 允许 nt-level insert/delete/substitute；CDS 默认只允许 synonymous codon substitution。",
            "CTMC edit rates：模型预测各区域插入/删除/替换强度，采样轨迹就是可解释 edit script。",
            "硬约束：legal sequence、CDS frame、protein identity、budget、motif/frame controls 作为 exact-1 gate。",
            "优化方式：候选生成 + proposal ranker + multi-objective teacher + source-aware cascade。",
        ],
    )
    slides_meta.append(("算法原理", ["CTMC edit rates + region grammar + hard constraints"]))

    # 11 architecture
    s = slide(prs, "模型架构：mRNA-native 组件")
    add_table(
        s,
        ["模块", "作用", "对应生物约束"],
        [
            ["Region FiLM", "按 5'UTR/CDS/3'UTR 调制 hidden states", "区域功能不同"],
            ["Codon phase features", "CDS 内记录 phase 0/1/2", "reading frame"],
            ["Synonymous CDS mask", "非同义替换概率置零", "protein identity exact-1"],
            ["RoPE / long context", "支持 full transcript context", "跨区域 interaction"],
            ["Aux structure proxy", "学习结构/可及性 proxy", "起始区和稳定性"],
            ["Proposal ranker", "重排候选，蒸馏 oracle 偏好", "top-1 design quality"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("模型架构", ["Region FiLM / phase / synonymous mask / ranker"]))

    # 12 training
    s = slide(prs, "训练策略：Stage A + ranker + downstream audits")
    add_metric_cards(
        s,
        [
            ("A100-max", "117.26M", "trainable params"),
            ("语料", "54,680", "GENCODE cleaned records"),
            ("目标步数", "100k", "Stage A max per seed"),
            ("当前布局", "4+2 GPUs", "4 A100-max + 2 MIG tiny"),
        ],
    )
    add_bullets(
        s,
        [
            "A100-max：GPU0/2/4/5 跑 seed0/1/2/5，100k steps，max UTR context 参考 mRNA-LM。",
            "MIG tiny：GPU6/7 仅 4.75GB，可跑 2.04M tiny config，不跑 117M max config。",
            "seed0 post-eval watcher：训练完成后自动接 proposal-ranking、ranker distill、T5 multiseed、compare、spectrum audit。",
            "P3/P4 controlled sweep：data256/1024 × tiny/small × steps200/500 已排队，用于真正 scale-law audit。",
        ],
        y=2.85,
        size=14,
    )
    slides_meta.append(("训练策略", ["A100-max + MIG tiny + post-eval queue"]))

    # 13 tasks
    s = slide(prs, "下游任务：T1-T7 全链路")
    add_table(
        s,
        ["任务", "核心问题", "关键指标"],
        [
            ["T1 Validity/Oracle", "生成是否合法且 proxy 是否提升", "legal, oracle TE/MRL"],
            ["T2 Distribution", "TE 提升是否破坏分布", "k-mer JS, codon KL, GC/length"],
            ["T3 Diversity/Novelty", "是否复制训练集", "novelty, exact match, diversity"],
            ["T4 Protein/CDS", "CDS 是否保持蛋白并优化 codon", "protein=1, CAI/GC, codon recovery"],
            ["T5 Edit Budget", "小预算下优化 TE", "within budget, edit distance, TE delta"],
            ["T6 Length Control", "长度是否可控", "abs length error, TE side effect"],
            ["T7 Motif/Frame", "motif 插删与 frame 安全", "success, frame/protein/budget"],
        ],
        h=5.0,
        font=8,
    )
    slides_meta.append(("下游任务", ["T1-T7"]))

    # 14 evaluation
    s = slide(prs, "测评方法：统计、约束和 claim gates")
    add_bullets(
        s,
        [
            "默认 10 seeds；mean ± bootstrap 95% CI；同一输入 paired test / Wilcoxon；必要时多重校正。",
            "硬约束 exact-1：legal、protein identity、reading frame、within budget、motif/frame controls。",
            "优化 oracle 与评测 oracle 分离；内部 TE/结构头只做引导，不作为最终主评测。",
            "positive SOTA claim gate：external real metrics、full de novo、real MPRA/stability、true scale-law、wet-lab 全部未完成时必须关闭。",
            f"当前 readiness：all_ready_for_sota_claim_audit={readiness.get('all_ready_for_sota_claim_audit')}; positive_sota_claim_ready={readiness.get('positive_sota_claim_ready')}",
        ],
    )
    slides_meta.append(("测评方法", ["10 seeds + paired tests + exact-1 gates"]))

    # 15 T5 head256 chart
    s = slide(prs, "当前主结果：head256 多目标 ranker 显著提升")
    labels = ["TE-only", "Pareto", "Scalar", "GRPO"]
    values = [t5_head256["te_only"] or 0, t5_head256["pareto"] or 0, t5_head256["scalar"] or 0, t5_head256["grpo"] or 0]
    add_bar_chart(s, labels, values, title="T5 delta_oracle_te_vs_source (head256, 10 seeds)")
    add_bullets(
        s,
        [
            "GRPO/scalar/pareto 三种 multi-objective fusion 均显著优于同 recipe TE-only 与旧 hardneg_v2。",
            "最佳 GRPO delta TE = +0.01114；TE-only = +0.00348；所有 hard constraints = 1.0。",
            "三种 fusion 彼此不显著：不声称单一 fusion 绝对最优。",
        ],
        y=5.45,
        size=12,
    )
    slides_meta.append(("head256 主结果", ["MO fusion 显著提升"]))

    # 16 head1024
    s = slide(prs, "scale-up 结果：head1024 收益收缩，必须诚实表述")
    labels = ["hardneg v2", "TE-only", "Pareto"]
    values = [t5_head1024["hardneg"] or 0, t5_head1024["te_only"] or 0, t5_head1024["pareto"] or 0]
    add_bar_chart(s, labels, values, title="T5 delta_oracle_te_vs_source (head1024, 10 seeds)", color=COLORS["purple"])
    add_bullets(
        s,
        [
            "head1024 上 TE-only 自身变强：+0.00846；Pareto = +0.00927。",
            "Pareto vs TE-only 仅 +0.00081，paired p=0.05047：borderline/trend，不写严格显著。",
            "三种 fusion 仍显著优于旧 hardneg_v2；但 head256 的强阳性未同等复现在更强对照 head1024 上。",
        ],
        y=5.45,
        size=12,
    )
    slides_meta.append(("head1024 scale-up", ["收益收缩，borderline"]))

    # 17 distribution
    s = slide(prs, "分布与序列谱审计：避免只骗低阶 oracle")
    add_metric_cards(
        s,
        [
            ("Base comp L1", fmt(spectrum.get("base_composition_full_l1"), 4), "candidate vs source"),
            ("Length Δ", fmt(spectrum.get("length_mean_delta"), 4), "mean length delta"),
            ("GC Δ", fmt(spectrum.get("gc_mean_delta"), 4), "mean GC delta"),
            ("k-mer L1", fmt(spectrum.get("kmer_l1"), 4), "multi-scale spectrum"),
        ],
    )
    add_bullets(
        s,
        [
            "已补 A/C/G/U base composition、region-wise composition、length/GC histogram、k-mer/codon-pair spectrum。",
            "当前 head32 ranker-full1k audit ready_for_distribution_figure_audit=True；但不是 wet-lab 或 de novo SOTA 证据。",
            "下一步：把 spectrum audit 和 TE delta / motif / codon-pair bias 联合成 paper figure。",
        ],
        y=2.8,
        size=14,
    )
    slides_meta.append(("分布审计", ["多尺度序列谱"]))

    # 18 T4
    s = slide(prs, "T4 protein-conditioned CDS：head1024 已跑通")
    add_metric_cards(
        s,
        [
            ("Targets", "1024", "protein-CDS pairs"),
            ("Protein identity", "1.0", "exact hard gate"),
            ("Designed CAI", fmt(t4_cds.get("mean_designed_cai"), 4), "native 0.6937"),
            ("Codon recovery", fmt(t4_metrics.get("mean_native_codon_recovery"), 4), "native codon level"),
        ],
    )
    add_bullets(
        s,
        [
            f"Protein-conditioned design: CAI delta vs native = {fmt(t4_cds.get('mean_designed_vs_native_cai_delta'), 4, signed=True)}；mean codon changes = {fmt(t4_cds.get('mean_codon_changes'), 2)}。",
            f"Codon metrics: synonymous substitution fraction = {fmt(t4_metrics.get('mean_native_synonymous_substitution_fraction'), 4)}；nonsynonymous = {fmt(t4_metrics.get('mean_native_nonsynonymous_substitution_fraction'), 4)}。",
            f"Codon usage KL = {fmt(t4_metrics.get('designed_vs_native_codon_usage_kl'), 4)}；codon-pair KL = {fmt(t4_metrics.get('designed_vs_native_codon_pair_kl'), 4)}。",
            "解释：当前 DP 是强 CAI/GC3 同义重编码，不是 native codon reconstruction；不能写成 codonGPT/Prot2RNA 外部复现。",
        ],
        y=2.8,
        size=13,
    )
    slides_meta.append(("T4 结果", ["head1024 protein-conditioned CDS/codon metrics"]))

    # 19 GC sweep
    s = slide(prs, "T4 CAI-GC Pareto：高 CAI 与 GC 控制存在张力")
    front = t4_sweep.get("pareto_front", []) if isinstance(t4_sweep, dict) else []
    rows = []
    for point in front[:5]:
        summ = point.get("summary", {})
        rows.append([
            str(point.get("gc_weight")),
            fmt(summ.get("mean_designed_cai"), 4),
            fmt(summ.get("mean_designed_gc"), 4),
            fmt(summ.get("mean_abs_gc_error"), 4),
            fmt(summ.get("protein_identity_eq_1_fraction"), 1),
        ])
    add_table(s, ["gc_weight", "CAI", "GC", "|GC-target|", "identity"], rows, h=2.9, font=9)
    add_bullets(
        s,
        [
            f"GC sweep audit ready={t4_sweep_audit.get('ready_for_pareto_claim_audit')}；8 weights × 1024 targets；all-point identity exact-1。",
            "GC-target best：gc_weight 8/16，CAI≈0.8680，GC≈0.5951，abs GC error≈0.0454。",
            "CAI-max：CAI≈1.0，但 GC 更高；这提示 codon optimization 需要 Pareto report，不能只报 CAI。",
        ],
        y=4.55,
        size=12,
    )
    slides_meta.append(("CAI-GC Pareto", ["T4 tradeoff"]))

    # 20 T6/T7
    s = slide(prs, "T6/T7：长度控制与 motif/frame 构造性能力")
    add_table(
        s,
        ["任务", "当前证据", "诚实解释"],
        [
            ["T6 head256/head1024", "5 个 target deltas 均 complete；hard constraints exact-1", "长度控制安全；正向加长 +15/+30 会损伤 proxy TE"],
            ["T7 motif insert/excise", "success≈0.993；frame/protein/budget=1.0", "证明构造性 motif 编辑；不证明 motif 策略已优化 TE"],
            ["Edit budget curve", "budget 1/2/3/5/10 全曲线", "TE delta 随预算上升；仍需 wet-lab/外部 baseline"],
        ],
        h=3.2,
        font=10,
    )
    add_bullets(
        s,
        [
            "T6 是 MEF 的强差异化：变长编辑天然支持 length-conditioned control。",
            "T7 是可解释轨迹优势：可以直接审计 motif 插入/切除和 reading frame。",
        ],
        y=4.7,
        size=14,
    )
    slides_meta.append(("T6/T7", ["长度和 motif/frame"]))

    # 21 External baselines
    s = slide(prs, "外部 baseline readiness：协议就绪，但 executable 未配置")
    table3 = load_json("docs/paper_table3_external_baseline_readiness.json")
    rows = []
    for row in table3.get("rows", [])[:4]:
        rows.append([
            str(row.get("model_name", row.get("candidate", ""))),
            str(row.get("status", "")),
            str(row.get("expected_outputs", row.get("notes", "")))[:80],
        ])
    add_table(s, ["模型", "状态", "说明"], rows, h=3.4, font=8)
    add_bullets(
        s,
        [
            "LinearDesign / EnsembleDesign / codonGPT / UTailoR 当前均为 not_configured。",
            "已经建立 dry-run registry、dataset hash、split、seed、runtime、hardware contract。",
            "外部 SOTA metric claim 必须等 executable_ready 且真实 adapter 写入 measured metrics。",
        ],
        y=4.95,
        size=13,
    )
    slides_meta.append(("外部 baseline", ["协议 ready；未配置 executable"]))

    # 22 Data readiness
    s = slide(prs, "数据 scale-up readiness：当前缺口")
    add_metric_cards(
        s,
        [
            ("GENCODE", "Ready", "manifest + records SHA"),
            ("RefSeq", str(data_ready.get("refseq_build_status")), "records/manifest missing"),
            ("MPRA/TE", "Not ready", "real table missing"),
            ("Stability", "Not ready", "real labels missing"),
        ],
    )
    add_bullets(
        s,
        [
            "family_split_protocol_ready=True，但 family_leakage_ready=False：RefSeq cross-corpus leakage 未完成。",
            "dataset_manifest_contract_ready=False：RefSeq、MPRA、stability manifests 仍缺。",
            "eFold-inspired gate 已加入文档：data scale-up 后续必须报告 diversity/complexity/cross-source panels。",
        ],
        y=2.85,
        size=14,
    )
    slides_meta.append(("数据 readiness", ["GENCODE ready；其它缺口"]))

    # 23 Current remote execution
    s = slide(prs, "当前远端运行状态")
    add_metric_cards(
        s,
        [
            ("Artifacts", f"{remote.get('n_artifacts_present', 0)} present", f"{remote.get('n_artifacts_missing', 0)} missing"),
            ("Dynamic procs", str(remote.get("n_dynamic_processes", "NA")), "remote status"),
            ("T6 head1024", "5/5 complete", "length deltas"),
            ("Readiness", str(remote.get("readiness_present")), "status snapshot"),
        ],
    )
    add_bullets(
        s,
        [
            "A100-max：4 seeds 继续训练；seed0 post-eval watcher 等待 Stage A 完成后自动衔接 downstream。",
            "MIG tiny：GPU6/7 跑 tiny config，利用 4.75GB MIG 分区。",
            "RefSeq build：仍 queued_or_running，raw/records/manifest 未落盘。",
            "head1024 protein-conditioned T4 已 complete 并被 remote status 采集。",
        ],
        y=2.85,
        size=14,
    )
    slides_meta.append(("远端状态", ["A100/MIG/RefSeq/post-eval"]))

    # 24 Model capability
    s = slide(prs, "当前模型能力：能做什么，不能 claim 什么")
    add_table(
        s,
        ["能力", "当前状态", "边界"],
        [
            ["Full transcript local optimization", "可运行；T1-T7 proxy evidence ready", "不是 full de novo SOTA"],
            ["Hard constraints", "legal/protein/frame/budget exact-1", "必须持续 gate"],
            ["TE proxy uplift", "head256 严格正向；head1024 收益收缩", "不是 wet-lab TE"],
            ["Protein-conditioned CDS", "head1024 ready", "不是 codonGPT/LinearDesign 复现"],
            ["Distribution audit", "base/region/k-mer/codon-pair ready", "不是天然功能证明"],
            ["Scale-law", "partial axes ready，controlled sweep running", "true scale-law claim false"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("模型能力边界", ["当前能做和不能 claim"]))

    # 25 Claim gate
    s = slide(prs, "SOTA / 论文 claim gate")
    add_bullets(
        s,
        [
            f"all_ready_for_sota_claim_audit = {readiness.get('all_ready_for_sota_claim_audit')}",
            f"positive_sota_claim_ready = {readiness.get('positive_sota_claim_ready')}",
            "允许写：constrained local full-length mRNA optimization/reranking with proxy/offline T1-T7 evidence。",
            "不能写：full de novo SOTA、wet-lab expression superiority、external SOTA reproduction、true scale-law、real MPRA/stability claim。",
            "阻断原因：external real metrics missing；full de novo evidence missing；real MPRA/stability missing；true scale-law missing；head1024 vs TE-only not strict；wet-lab missing。",
        ],
    )
    slides_meta.append(("Claim gate", ["positive SOTA claim 关闭"]))

    # 26 Roadmap
    s = slide(prs, "未来规划：P0-P4")
    add_table(
        s,
        ["优先级", "任务", "交付物"],
        [
            ["P0", "完成 A100-max Stage A + post-eval", "seed0/1/2/5 checkpoints + T1-T7 + compare + spectrum"],
            ["P0", "外部 executable baseline", "LinearDesign/EnsembleDesign/codonGPT/UTailoR measured metrics"],
            ["P1", "真实 MPRA/stability 数据", "official split + manifest + held-out predictor audits"],
            ["P1", "Frozen FM embedding cache", "mRNA-LM/Helix/CodonFM leakage-free adapter probe"],
            ["P2", "eFold-inspired cross-family panels", "diversity profile + complexity bucket + cross-source gap"],
            ["P3/P4", "true data/model/step scale-law", "controlled sweep + downstream trend audit"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("未来规划", ["P0-P4"]))

    # 27 Leadership asks
    s = slide(prs, "需要领导支持的资源与决策")
    add_bullets(
        s,
        [
            "计算资源：A100 训练周期较长，建议保留 4 张完整 A100 跑 max seeds，并让 MIG 分区继续跑 tiny/aux evidence。",
            "数据授权：MPRA/TE、half-life/stability、外部 FM 权重/许可需要明确下载与使用边界。",
            "外部 baseline：需要配置 LinearDesign/EnsembleDesign/codonGPT/UTailoR executable 或容器。",
            "实验验证：若目标是领导层/论文级正向 SOTA，后续必须规划 wet-lab 或至少独立真实 predictor validation。",
            "汇报口径：当前应主打“硬约束安全 + 代理评测完整 + scale-up 诚实审计”，避免过早 SOTA 口号。",
        ],
    )
    slides_meta.append(("资源诉求", ["算力、数据、baseline、wet-lab"]))

    # 28 Appendix: artifacts
    s = slide(prs, "附录：关键 artifact 路径")
    add_table(
        s,
        ["类别", "路径"],
        [
            ["Roadmap", "docs/next_steps_sota_roadmap.md"],
            ["任务协议", "docs/downstream_tasks.md"],
            ["数据调研", "docs/mrna_dataset_survey.md"],
            ["SOTA readiness", "docs/sota_readiness_audit_head256.json"],
            ["Remote status", "docs/remote_execution_status.json"],
            ["Head1024 T4", "benchmark/t4_protein_identity_cai_gc_report_head1024.json"],
            ["Data readiness", "docs/data_scaleup_readiness.json"],
            ["External baselines", "docs/paper_table3_external_baseline_readiness.json"],
        ],
        h=4.6,
        font=9,
    )
    slides_meta.append(("附录路径", ["关键 artifact"]))

    # 29 Appendix: exact numbers
    s = slide(prs, "附录：关键数字汇总")
    add_table(
        s,
        ["指标", "数值", "备注"],
        [
            ["head256 GRPO delta TE", fmt(t5_head256["grpo"], 5, signed=True), "vs source; 10 seeds"],
            ["head1024 Pareto delta TE", fmt(t5_head1024["pareto"], 5, signed=True), "borderline vs TE-only"],
            ["head1024 TE-only delta TE", fmt(t5_head1024["te_only"], 5, signed=True), "strong control"],
            ["T4 head1024 designed CAI", fmt(t4_cds.get("mean_designed_cai"), 5), "protein exact-1"],
            ["T4 head1024 codon recovery", fmt(t4_metrics.get("mean_native_codon_recovery"), 5), "native codon level"],
            ["Spectrum base L1", fmt(spectrum.get("base_composition_full_l1"), 5), "distribution audit"],
            ["GENCODE records", "54,680", "cleaned full transcripts"],
            ["A100-max params", "117.26M", "trainable"],
        ],
        h=4.8,
        font=9,
    )
    slides_meta.append(("附录数字", ["关键指标"]))

    # 30 closing
    s = slide(prs, "总结：当前最适合的汇报口径")
    add_bullets(
        s,
        [
            "我们已经建立了 mRNA-native constrained edit-flow 的完整工程闭环：数据、训练、候选、ranker、T1-T7、readiness gate。",
            "当前最强证据是内部 proxy/offline：head256 多目标显著，head1024 趋势收缩但约束稳定，T4 protein-conditioned head1024 完整。",
            "项目真正的论文壁垒不是“模型更大”，而是：全长区域编辑文法 + 硬约束 exact-1 + cross-family/data-diversity 泛化治理。",
            "下一阶段只要补齐真实数据、外部 executable 和 frozen FM，就能把内部 evidence 转成可对外竞争的 paper-grade 结果。",
        ],
        y=1.35,
        size=17,
        color=COLORS["navy"],
    )
    slides_meta.append(("总结", ["硬约束安全、proxy 完整、下一步补外部证据"]))

    for i, sl in enumerate(prs.slides, start=1):
        if i != 1:
            add_footer(sl, i)

    prs.save(OUT_PPTX)
    with open(OUT_MD, "w", encoding="utf-8") as fh:
        fh.write("# mRNA-EditFlow Leadership Briefing Outline\n\n")
        fh.write(f"- PPTX: `{OUT_PPTX}`\n")
        fh.write(f"- Generated: `{date.today().isoformat()}`\n")
        fh.write(f"- Slides: `{len(slides_meta)}`\n\n")
        for i, (title, bullets) in enumerate(slides_meta, start=1):
            fh.write(f"## {i}. {title}\n\n")
            for b in bullets:
                fh.write(f"- {b}\n")
            fh.write("\n")
    return OUT_PPTX, OUT_MD, len(slides_meta)


if __name__ == "__main__":
    pptx, md, n = build_deck()
    print(json.dumps({"pptx": str(pptx), "outline": str(md), "slides": n}, ensure_ascii=False, sort_keys=True))
